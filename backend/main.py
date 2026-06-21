import uvicorn
import os
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import tempfile
import boto3
import jwt
import bcrypt
from datetime import datetime, timedelta, timezone
import secrets
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from fastapi import Depends
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
import base64

from dotenv import load_dotenv

from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, CSVLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from pptx import Presentation
import openpyxl

load_dotenv()

app = FastAPI(title="Chat with PDF using Gemini")

# CORS setup for React frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # In production, change this to specific domains
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

SECRET_KEY = os.getenv("JWT_SECRET", "supersecretkey")
ALGORITHM = "HS256"

# DynamoDB setup
dynamodb = boto3.resource('dynamodb', region_name=os.getenv("AWS_REGION", "us-east-1"))
USERS_TABLE_NAME = os.getenv("DYNAMODB_TABLE", "ChatBot_User")

def get_users_table():
    return dynamodb.Table(USERS_TABLE_NAME)

# S3 setup
s3 = boto3.client('s3', region_name=os.getenv("AWS_REGION", "us-east-1"))
S3_BUCKET = os.getenv("S3_BUCKET_NAME", "")

security = HTTPBearer()

def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    token = credentials.credentials
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise HTTPException(status_code=401, detail="Invalid authentication credentials")
        return email
    except jwt.PyJWTError:
        raise HTTPException(status_code=401, detail="Invalid authentication credentials")

# Global store for the vector database (in-memory for demo purposes)
# For production, consider persistent storage like ChromaDB or Postgres pgvector
app.state.vector_store = None
app.state.active_file = None

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".txt", ".csv", ".xlsx",
    ".png", ".jpg", ".jpeg", ".webp"
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}

# Set up Gemini clients
# Make sure GOOGLE_API_KEY is set in .env
try:
    # Initialize embeddings
    embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")

    # Initialize LLM
    llm = ChatGoogleGenerativeAI(model="gemini-3.1-flash-lite-preview")
except Exception as e:
    print(f"Error initializing Gemini clients: {e}")
    embeddings = None
    llm = None

class ChatRequest(BaseModel):
    query: str

class RegisterRequest(BaseModel):
    username: str
    email: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str

class ForgotPasswordRequest(BaseModel):
    email: str

class ResetPasswordRequest(BaseModel):
    email: str
    token: str
    new_password: str

class VerifyEmailRequest(BaseModel):
    email: str
    code: str

def send_verification_email(to_email: str, code: str) -> bool:
    smtp_user = os.getenv("SMTP_USER")
    smtp_pass = os.getenv("SMTP_PASSWORD")
    smtp_host = os.getenv("SMTP_HOST", "smtp.gmail.com")
    smtp_port_str = os.getenv("SMTP_PORT", "587")

    if not smtp_user or not smtp_pass:
        print("Warning: SMTP not configured. Skipping verification email.")
        return False

    try:
        smtp_port = int(smtp_port_str)
        msg = MIMEMultipart()
        msg['From'] = smtp_user
        msg['To'] = to_email
        msg['Subject'] = "Email Verification Code - Insight AI"

        body = (
            "<html><body style='font-family: Arial, sans-serif; line-height: 1.6; color: #333333;'>"
            "<div style='max-width: 600px; margin: 20px auto; padding: 30px; border: 1px solid #e2e8f0; border-radius: 12px; background-color: #ffffff;'>"
            "<div style='text-align: center; margin-bottom: 24px;'>"
            "<h2 style='color: #4f46e5; margin: 0; font-size: 1.5rem; font-weight: 700;'>Insight AI</h2>"
            "<p style='font-size: 0.875rem; color: #64748b;'>Email Verification</p></div>"
            "<p>Hello,</p>"
            "<p>Thanks for signing up! Use the 6-digit code below to verify your email. This code expires in 15 minutes.</p>"
            "<div style='text-align: center; margin: 32px 0;'>"
            f"<div style='display: inline-block; font-size: 32px; font-weight: 700; letter-spacing: 6px; color: #4f46e5; padding: 16px 32px; border: 2px dashed #e2e8f0; border-radius: 8px; background-color: #f8fafc; font-family: monospace;'>{code}</div>"
            "</div>"
            "<p>If you did not create an account, you can safely ignore this email.</p>"
            "<hr style='border: 0; border-top: 1px solid #e2e8f0; margin: 28px 0;' />"
            "<p style='font-size: 11px; color: #94a3b8; text-align: center;'>This is an automated message. Please do not reply.</p>"
            "</div></body></html>"
        )
        msg.attach(MIMEText(body, 'html'))
        server = smtplib.SMTP(smtp_host, smtp_port)
        server.starttls()
        server.login(smtp_user, smtp_pass)
        server.sendmail(smtp_user, to_email, msg.as_string())
        server.quit()
        print(f"Verification email sent to {to_email}")
        return True
    except Exception as e:
        print(f"Error sending verification email: {e}")
        return False

@app.post("/register")
def register(user: RegisterRequest):
    table = get_users_table()

    response = table.get_item(Key={'email': user.email})
    if 'Item' in response:
        existing = response['Item']
        if existing.get('email_verified', False):
            raise HTTPException(status_code=400, detail="Email already registered")
        # Allow re-registration if not yet verified (resend code)

    hashed_password = bcrypt.hashpw(user.password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
    code = f"{secrets.randbelow(900000) + 100000}"
    expiry = int((datetime.now(timezone.utc) + timedelta(minutes=15)).timestamp())

    table.put_item(Item={
        'email': user.email,
        'username': user.username,
        'password': hashed_password,
        'email_verified': False,
        'verification_code': code,
        'verification_code_expiry': expiry
    })

    email_sent = send_verification_email(user.email, code)

    if email_sent:
        return {"message": "Verification code sent to your email", "simulation": False}
    else:
        return {
            "message": "Account created (Simulation mode: SMTP not configured)",
            "code": code,
            "simulation": True
        }

@app.post("/verify-email")
def verify_email(req: VerifyEmailRequest):
    table = get_users_table()
    response = table.get_item(Key={'email': req.email})
    if 'Item' not in response:
        raise HTTPException(status_code=404, detail="Email not found")

    item = response['Item']
    stored_code = item.get('verification_code')
    stored_expiry = item.get('verification_code_expiry')

    if not stored_code or stored_code != req.code:
        raise HTTPException(status_code=400, detail="Invalid verification code")

    current_time = int(datetime.now(timezone.utc).timestamp())
    if stored_expiry and current_time > int(stored_expiry):
        raise HTTPException(status_code=400, detail="Verification code has expired")

    table.update_item(
        Key={'email': req.email},
        UpdateExpression="SET email_verified = :v REMOVE verification_code, verification_code_expiry",
        ExpressionAttributeValues={':v': True}
    )
    return {"message": "Email verified successfully"}

@app.post("/login")
def login(user: LoginRequest):
    table = get_users_table()
    response = table.get_item(Key={'email': user.email})

    if 'Item' not in response:
        raise HTTPException(status_code=400, detail="Incorrect email or password")

    item = response['Item']
    if not item.get('email_verified', False):
        raise HTTPException(status_code=403, detail="Please verify your email before logging in")

    stored_password = item['password']
    if not bcrypt.checkpw(user.password.encode('utf-8'), stored_password.encode('utf-8')):
        raise HTTPException(status_code=400, detail="Incorrect email or password")

    username = item['username']
    token = jwt.encode(
        {"sub": user.email, "exp": datetime.now(timezone.utc) + timedelta(hours=24)},
        SECRET_KEY,
        algorithm=ALGORITHM
    )
    return {"access_token": token, "token_type": "bearer", "username": username}

def send_reset_email(to_email: str, token: str) -> bool:
    smtp_user = os.getenv("SMTP_USER")
    smtp_pass = os.getenv("SMTP_PASSWORD")
    smtp_host = os.getenv("SMTP_HOST", "smtp.gmail.com")
    smtp_port_str = os.getenv("SMTP_PORT", "587")
    
    if not smtp_user or not smtp_pass:
        print("Warning: SMTP_USER or SMTP_PASSWORD not configured. Skipping email send.")
        return False
        
    try:
        smtp_port = int(smtp_port_str)
        # Create message container
        msg = MIMEMultipart()
        msg['From'] = smtp_user
        msg['To'] = to_email
        msg['Subject'] = "Password Reset Verification Code - Insight AI"
        
        body = (
            "<html><body style='font-family: Arial, sans-serif; line-height: 1.6; color: #333333;'>"
            "<div style='max-width: 600px; margin: 20px auto; padding: 30px; border: 1px solid #e2e8f0; border-radius: 12px; background-color: #ffffff;'>"
            "<div style='text-align: center; margin-bottom: 24px;'>"
            "<h2 style='color: #4f46e5; margin: 0; font-size: 1.5rem; font-weight: 700;'>Insight AI</h2>"
            "<p style='font-size: 0.875rem; color: #64748b;'>Secure Password Reset</p></div>"
            "<p>Hello,</p>"
            "<p>We received a request to reset your password. Use the 6-digit code below to proceed. This code is valid for 15 minutes.</p>"
            "<div style='text-align: center; margin: 32px 0;'>"
            f"<div style='display: inline-block; font-size: 32px; font-weight: 700; letter-spacing: 6px; color: #4f46e5; padding: 16px 32px; border: 2px dashed #e2e8f0; border-radius: 8px; background-color: #f8fafc; font-family: monospace;'>{token}</div>"
            "</div>"
            "<p>If you did not make this request, you can safely ignore this email.</p>"
            "<hr style='border: 0; border-top: 1px solid #e2e8f0; margin: 28px 0;' />"
            "<p style='font-size: 11px; color: #94a3b8; text-align: center;'>This is an automated message. Please do not reply.</p>"
            "</div></body></html>"
        )
        msg.attach(MIMEText(body, 'html'))
        
        # Connect to Gmail SMTP server using TLS
        server = smtplib.SMTP(smtp_host, smtp_port)
        server.starttls()
        server.login(smtp_user, smtp_pass)
        server.sendmail(smtp_user, to_email, msg.as_string())
        server.quit()
        print(f"Password reset email sent to {to_email}")
        return True
    except Exception as e:
        print(f"Error sending email: {e}")
        return False

@app.post("/forgot-password")
def forgot_password(req: ForgotPasswordRequest):
    table = get_users_table()
    response = table.get_item(Key={'email': req.email})
    if 'Item' not in response:
        raise HTTPException(status_code=404, detail="Email not registered")

    if not response['Item'].get('email_verified', False):
        raise HTTPException(status_code=403, detail="Please verify your email first before resetting password")
    
    # Generate 6-digit PIN/token
    token = f"{secrets.randbelow(900000) + 100000}"
    # Expire in 15 minutes
    expiry = int((datetime.now(timezone.utc) + timedelta(minutes=15)).timestamp())
    
    table.update_item(
        Key={'email': req.email},
        UpdateExpression="SET reset_token = :token, reset_token_expiry = :expiry",
        ExpressionAttributeValues={
            ':token': token,
            ':expiry': expiry
        }
    )
    
    # Try sending email
    email_sent = send_reset_email(req.email, token)
    
    if email_sent:
        return {"message": "Verification code sent to your email", "simulation": False}
    else:
        return {
            "message": "Reset code generated (Simulation mode: SMTP not configured)",
            "token": token,
            "simulation": True
        }

@app.post("/reset-password")
def reset_password(req: ResetPasswordRequest):
    table = get_users_table()
    response = table.get_item(Key={'email': req.email})
    if 'Item' not in response:
        raise HTTPException(status_code=404, detail="Email not registered")
    
    item = response['Item']
    stored_token = item.get('reset_token')
    stored_expiry = item.get('reset_token_expiry')
    
    if not stored_token or stored_token != req.token:
        raise HTTPException(status_code=400, detail="Invalid verification code")
        
    current_time = int(datetime.now(timezone.utc).timestamp())
    if stored_expiry and current_time > int(stored_expiry):
        raise HTTPException(status_code=400, detail="Verification code has expired")
        
    hashed_password = bcrypt.hashpw(req.new_password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
    
    table.update_item(
        Key={'email': req.email},
        UpdateExpression="SET password = :p, email_verified = :v REMOVE reset_token, reset_token_expiry",
        ExpressionAttributeValues={
            ':p': hashed_password,
            ':v': True
        }
    )
    return {"message": "Password reset successfully"}

def process_file_to_vectorstore(tmp_file_path: str, filename: str):
    ext = os.path.splitext(filename)[1].lower()

    if ext in IMAGE_EXTENSIONS:
        # Use Gemini Vision to extract text from image
        with open(tmp_file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")
        mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".webp": "image/webp"}
        vision_llm = ChatGoogleGenerativeAI(model="gemini-3.1-flash-lite-preview")
        from langchain_core.messages import HumanMessage
        msg = HumanMessage(content=[
            {"type": "text", "text": "Extract and return all text content visible in this image in detail."},
            {"type": "image_url", "image_url": {"url": f"data:{mime_map[ext]};base64,{image_data}"}}
        ])
        extracted_text = vision_llm.invoke([msg]).content
        docs = [Document(page_content=extracted_text, metadata={"source": filename})]
    elif ext == ".pdf":
        docs = PyPDFLoader(tmp_file_path).load()
    elif ext == ".docx":
        docs = Docx2txtLoader(tmp_file_path).load()
    elif ext == ".pptx":
        prs = Presentation(tmp_file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        docs = [Document(page_content="\n".join(text_parts), metadata={"source": filename})]
    elif ext == ".txt":
        with open(tmp_file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        docs = [Document(page_content=text, metadata={"source": filename})]
    elif ext == ".csv":
        docs = CSVLoader(tmp_file_path).load()
    elif ext == ".xlsx":
        wb = openpyxl.load_workbook(tmp_file_path, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        docs = [Document(page_content="\n".join(text_parts), metadata={"source": filename})]
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    splits = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(docs)
    app.state.vector_store = FAISS.from_documents(splits, embeddings)
    return len(splits)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...), current_user: str = Depends(get_current_user)):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(SUPPORTED_EXTENSIONS)}")
    if embeddings is None or llm is None:
        raise HTTPException(status_code=500, detail="Gemini clients not initialized properly.")
    if not S3_BUCKET:
        raise HTTPException(status_code=500, detail="S3 bucket not configured.")
    try:
        content = await file.read()
        s3_key = f"{current_user}/{file.filename}"
        s3.put_object(Bucket=S3_BUCKET, Key=s3_key, Body=content)

        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            tmp_file.write(content)
            tmp_file_path = tmp_file.name

        chunks = process_file_to_vectorstore(tmp_file_path, file.filename)
        os.remove(tmp_file_path)
        app.state.active_file = file.filename

        return {"message": "File uploaded and processed successfully", "chunks": chunks, "filename": file.filename}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/files")
def list_files(current_user: str = Depends(get_current_user)):
    if not S3_BUCKET:
        raise HTTPException(status_code=500, detail="S3 bucket not configured.")
    try:
        prefix = f"{current_user}/"
        response = s3.list_objects_v2(Bucket=S3_BUCKET, Prefix=prefix)
        files = []
        for obj in response.get("Contents", []):
            key = obj["Key"]
            filename = key[len(prefix):]
            if filename:
                files.append({"filename": filename, "key": key, "size": obj["Size"], "last_modified": obj["LastModified"].isoformat()})
        return {"files": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error listing files: {str(e)}")

class SelectFileRequest(BaseModel):
    filename: str

@app.post("/files/select")
def select_file(req: SelectFileRequest, current_user: str = Depends(get_current_user)):
    if embeddings is None:
        raise HTTPException(status_code=500, detail="Gemini clients not initialized properly.")
    if not S3_BUCKET:
        raise HTTPException(status_code=500, detail="S3 bucket not configured.")
    try:
        ext = os.path.splitext(req.filename)[1].lower()
        s3_key = f"{current_user}/{req.filename}"
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            s3.download_fileobj(S3_BUCKET, s3_key, tmp_file)
            tmp_file_path = tmp_file.name
        chunks = process_file_to_vectorstore(tmp_file_path, req.filename)
        os.remove(tmp_file_path)
        app.state.active_file = req.filename
        return {"message": f"{req.filename} loaded successfully", "chunks": chunks, "filename": req.filename}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error loading file: {str(e)}")

@app.post("/chat")
async def chat_with_pdf(request: ChatRequest, current_user: str = Depends(get_current_user)):
    if app.state.vector_store is None:
        raise HTTPException(status_code=400, detail="Please upload a file first")
    
    if embeddings is None or llm is None:
        raise HTTPException(status_code=500, detail="Gemini clients not initialized properly.")
    try:
        retriever = app.state.vector_store.as_retriever(search_kwargs={"k": 3})
        system_prompt = (
            "You are an assistant for question-answering tasks. "
            "Use the following pieces of retrieved context to answer the question. "
            "If you don't know the answer, say that you don't know. "
            "Use three sentences maximum and keep the answer concise."
            "\n\n"
            "{context}"
        )
        prompt = ChatPromptTemplate.from_messages([
            ("system", system_prompt),
            ("human", "{input}"),
        ])
        def format_docs(docs):
            return "\n\n".join(doc.page_content for doc in docs)
        rag_chain = (
            {"context": retriever | format_docs, "input": RunnablePassthrough()}
            | prompt
            | llm
            | StrOutputParser()
        )

        response = rag_chain.invoke(request.query)
        return {"answer": response}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating answer: {str(e)}")
    
@app.get("/")
def read_root():
    return {"message": "Welcome to Chat with PDF GenAI API"}



if __name__ == "__main__":
    uvicorn.run("main:app", reload=True)